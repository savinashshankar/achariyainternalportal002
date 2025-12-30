"""
Content Extractor - Extract text from various file formats
"""

from typing import Dict, Optional
import io
import PyPDF2
from docx import Document
from google.cloud import storage
from core.drive_client import drive_client
from core.gemini_client import gemini_client


class ContentExtractor:
    """Extract content from various file formats"""
    
    def extract_from_file(self, file_id: str, mime_type: str) -> Dict:
        """
        Extract text content from file based on type
        
        Returns:
            {
                'text': str,
                'word_count': int,
                'page_count': int,
                'file_type': str
            }
        """
        
        if 'pdf' in mime_type:
            return self._extract_from_pdf(file_id)
        elif 'word' in mime_type or 'document' in mime_type:
            return self._extract_from_word(file_id)
        elif 'presentation' in mime_type:
            return self._extract_from_pptx(file_id)
        elif 'video' in mime_type:
            return self._extract_from_video(file_id)
        else:
            raise ValueError(f"Unsupported file type: {mime_type}")
    
    def _extract_from_pdf(self, file_id: str) -> Dict:
        """Extract text from PDF"""
        content = drive_client.download_file(file_id)
        
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
        text = ""
        
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        
        return {
            'text': text,
            'word_count': len(text.split()),
            'page_count': len(pdf_reader.pages),
            'file_type': 'pdf'
        }
    
    def _extract_from_word(self, file_id: str) -> Dict:
        """Extract text from Word document"""
        content = drive_client.download_file(file_id)
        
        doc = Document(io.BytesIO(content))
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        
        return {
            'text': text,
            'word_count': len(text.split()),
            'page_count': len(doc.paragraphs) // 20,  # Rough estimate
            'file_type': 'docx'
        }
    
    def _extract_from_pptx(self, file_id: str) -> Dict:
        """Extract text from PowerPoint"""
        from pptx import Presentation
        
        content = drive_client.download_file(file_id)
        prs = Presentation(io.BytesIO(content))
        
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return {
            'text': text,
            'word_count': len(text.split()),
            'page_count': len(prs.slides),
            'file_type': 'pptx'
        }
    
    def _extract_from_video(self, file_id: str) -> Dict:
        """Extract transcript from video using Gemini"""
        # For video, we'll use Gemini's multimodal capabilities
        # This is a placeholder - would need actual video processing
        
        # Download video
        content = drive_client.download_file(file_id)
        
        # Use Gemini to transcribe
        # Note: This requires Gemini multimodal API
        prompt = "Transcribe this educational video and provide a detailed text version."
        
        # Simplified for now - in production would use Gemini video API
        text = "Video transcription would go here using Gemini multimodal API"
        
        return {
            'text': text,
            'word_count': len(text.split()),
            'page_count': 1,
            'file_type': 'video'
        }


# Global instance
content_extractor = ContentExtractor()
